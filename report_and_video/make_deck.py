"""Build the 10-minute video presentation deck for the T8 term project.

Slide 1 is the compulsory IITG title template (rendered PNG, also the YouTube
thumbnail); slides 2-11 follow the submission guide's suggested flow. Palette
matches the template: navy + gold on white.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1B, 0x3A, 0x6B)
GOLD = RGBColor(0xC9, 0x9B, 0x1C)
INK = RGBColor(0x21, 0x21, 0x21)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0xED, 0xF1, 0xF7)   # pale navy tint for cards
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1B, 0x7A, 0x4B)
RED = RGBColor(0xA4, 0x2A, 0x2A)

# Paths are anchored to this script, not the working directory, so the deck
# rebuilds correctly from anywhere. Report figures live in the project root.
HERE = Path(__file__).resolve().parent
FIGURES = HERE.parent / "figures"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h, fill=None, line=None, shadow=False, round_=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if round_:
        shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill if fill else WHITE
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=15, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space_after=6, line_spacing=1.0):
    """runs: str, or list of paragraphs; each paragraph is str or list of
    (text, {bold, color, size, italic}) run tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for txt, fmt in para:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = "Calibri"
            f.size = Pt(fmt.get("size", size))
            f.bold = fmt.get("bold", bold)
            f.italic = fmt.get("italic", False)
            f.color.rgb = fmt.get("color", color)
    return tb


def header(s, kicker, title):
    text(s, Inches(0.6), Inches(0.38), Inches(12.1), Inches(0.3), kicker.upper(),
         size=13, color=GOLD, bold=True)
    text(s, Inches(0.6), Inches(0.68), Inches(12.1), Inches(0.75), title,
         size=32, color=NAVY, bold=True)


def bullets(s, x, y, w, h, items, size=15, space_after=10):
    """items: list of (lead, rest) -> bold lead + normal rest."""
    paras = []
    for lead, rest in items:
        runs = []
        if lead:
            runs.append((lead + "  ", {"bold": True, "color": NAVY}))
        runs.append((rest, {}))
        paras.append(runs)
    return text(s, x, y, w, h, paras, size=size, space_after=space_after,
                line_spacing=1.05)


def stat(s, x, y, w, value, label, color=NAVY, size=40):
    text(s, x, y, w, Inches(0.9), value, size=size, color=color, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, y + Inches(0.78), w, Inches(0.6), label, size=12.5, color=GREY,
         align=PP_ALIGN.CENTER)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


# ---------------------------------------------------------------- 1: title
s = slide()
s.shapes.add_picture(str(HERE / "title_slide.png"), 0, 0, SW, SH)
notes(s, "Hello, my name is Francois Schmitt, I am a student of the BSc Honours "
         "Data Science and AI Online Degree Programme at IIT Guwahati, and this "
         "is my Trimester 8 term project: Does Machine Learning Beat the Credit "
         "Scorecard? (state name, programme and project verbally, face visible)")

# ---------------------------------------------------------------- 2: problem
s = slide()
header(s, "The problem", "Probability of Default: the Highest-Stakes Model in Banking")
bullets(s, Inches(0.6), Inches(1.75), Inches(6.7), Inches(4.6), [
    ("The task:", "estimate the probability that a loan applicant defaults, "
     "using only information available on the day of application."),
    ("Why it matters:", "PD models drive accept/reject decisions, pricing, and, "
     "under Basel and IFRS 9, regulatory capital and provisions."),
    ("The status quo:", "banks still use logistic-regression scorecards, a "
     "forty-year-old technology, because they are stable and explainable."),
    ("The ML claim:", "benchmarking papers report boosted trees dominate, but "
     "they score models almost exclusively by AUC, the area under the ROC "
     "curve."),
    ("The question:", "does the ML advantage survive the criteria a lender "
     "actually needs?"),
], size=15.5)
x = Inches(7.7)
for i, (t1, t2) in enumerate([
        ("1. Discrimination", "Can the model rank bad borrowers above good ones? (area under the ROC curve, AUC; and KS)"),
        ("2. Calibration", "Can the predicted PD be taken at face value? (Brier, ECE, reliability)"),
        ("3. Profit", "Does the model make money at its accept/reject cutoff?")]):
    y = Inches(1.75) + i * Inches(1.55)
    box(s, x, y, Inches(5.0), Inches(1.3), fill=LIGHT, round_=True)
    text(s, x + Inches(0.3), y + Inches(0.17), Inches(4.4), Inches(0.4), t1,
         size=17, color=NAVY, bold=True)
    text(s, x + Inches(0.3), y + Inches(0.55), Inches(4.4), Inches(0.7), t2,
         size=13, color=INK)
notes(s, "The problem: PD estimation from origination-time data only. Banks use "
         "scorecards for regulatory reasons; the literature says ML wins but "
         "measures only discrimination. My project tests all three criteria.")

# ---------------------------------------------------------------- 3: tools
s = slide()
header(s, "Tools and technologies", "A Reproducible Python Stack, CPU Only")
cols = [
    ("Language", ["Python 3.13"]),
    ("Data", ["pandas", "NumPy", "kagglehub"]),
    ("Models", ["scikit-learn", "XGBoost", "LightGBM", "SciPy"]),
    ("Workflow", ["Jupyter", "VS Code", "Git / GitHub", "LaTeX"]),
]
for i, (title, items) in enumerate(cols):
    x = Inches(0.6) + i * Inches(3.2)
    box(s, x, Inches(1.9), Inches(2.9), Inches(3.4), fill=LIGHT, round_=True)
    text(s, x + Inches(0.3), Inches(2.15), Inches(2.3), Inches(0.4), title,
         size=18, color=GOLD, bold=True)
    text(s, x + Inches(0.3), Inches(2.65), Inches(2.3), Inches(2.4),
         items, size=15.5, color=INK, space_after=8)
text(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.8),
     [[("Everything is scripted:  ", {"bold": True, "color": NAVY}),
       ("one command reproduces the full study, about 18 minutes on a laptop; "
        "all seeds fixed, one shared vintage split.", {})]],
     size=15.5)
notes(s, "Quick tools mention: Python 3.13, pandas for data, scikit-learn plus "
         "XGBoost and LightGBM for models, everything in a modular package on "
         "GitHub; one command reproduces the study in about 18 minutes.")

# ---------------------------------------------------------------- 4: dataset
s = slide()
header(s, "Dataset", "Lending Club Personal Loans, Engineered Against Leakage")
stats = [
    ("621,022", "36-month loans, 2007-2015", 40),
    ("151 → 23", "columns after leakage whitelist", 40),
    ("12.6 → 14.9%", "default rate, train → test vintage", 30),
    ("30,000", "UCI Taiwan clients (robustness)", 40),
]
for i, (v, l, sz) in enumerate(stats):
    stat(s, Inches(0.5) + i * Inches(3.15), Inches(1.62), Inches(3.0), v, l, size=sz)

# what the loans actually are
box(s, Inches(0.6), Inches(2.95), Inches(12.1), Inches(1.95), fill=LIGHT, round_=True)
text(s, Inches(0.95), Inches(3.12), Inches(4.0), Inches(0.35),
     "Small unsecured consumer loans", size=15.5, color=NAVY, bold=True)
for i, (lead, rest) in enumerate([
        ("Size", "$500 to the $35,000 ceiling, median $10,000, $342 a month"),
        ("Price", "5.3% to 29.0%, mean near 12%"),
        ("Borrower", "median income $60,000, debt-to-income ratio 17%")]):
    y = Inches(3.55) + i * Inches(0.4)
    text(s, Inches(0.95), y, Inches(1.1), Inches(0.35), lead,
         size=12.5, color=GOLD, bold=True)
    text(s, Inches(2.1), y, Inches(4.4), Inches(0.35), rest, size=12.5, color=INK)
for i, (lead, rest) in enumerate([
        ("Purpose", "57% debt consolidation, 24% credit card: refinancing"),
        ("Housing", "46% mortgage, 43% renting, 10% owned"),
        ("Grades", "A to D is 96% of the book; E to G only 4.3%")]):
    y = Inches(3.55) + i * Inches(0.4)
    text(s, Inches(6.9), y, Inches(1.1), Inches(0.35), lead,
         size=12.5, color=GOLD, bold=True)
    text(s, Inches(8.05), y, Inches(4.4), Inches(0.35), rest, size=12.5, color=INK)
text(s, Inches(6.9), Inches(3.12), Inches(5.6), Inches(0.35),
     [[("FICO credit score:  ", {"bold": True, "color": NAVY}),
       ("floor 662, median 692, an approval cutoff", {"color": INK})]], size=12.5)

bullets(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(2.1), [
    ("The leakage filter is the hard work:", "of the 151 raw columns, most are "
     "post-origination servicing fields (payments received, recoveries, "
     "hardship flags) that trivially predict the outcome. I kept an explicit "
     "whitelist of 23 origination-time fields."),
    ("Out-of-time split:", "train on 2007-2013, calibrate on 2014, test on "
     "2015, exactly how a deployed model meets the future. The rising default "
     "rate across vintages is the drift a real model faces."),
], size=14.5, space_after=9)
notes(s, "Dataset: public Lending Club loans, and worth a moment on what they "
         "actually are. Small unsecured consumer loans, median ten thousand "
         "dollars over three years, mostly people refinancing more expensive "
         "credit-card debt. Note the FICO floor: from 2009 the minimum score "
         "in every single vintage is exactly 662, which is an approval cutoff, "
         "not a coincidence. Nobody genuinely subprime is in this sample, which "
         "is part of why the achievable AUCs sit near 0.68. Then the two key "
         "decisions: the origination-time whitelist against leakage, and the "
         "out-of-time vintage split.")

# ---------------------------------------------------------------- 5: approach
s = slide()
header(s, "Approach and code", "One Pipeline, Four Models, Identical Information")
box(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.95), fill=LIGHT, round_=True)
text(s, Inches(0.9), Inches(1.85), Inches(5.3), Inches(0.35),
     "The creditrisk package", size=15.5, color=NAVY, bold=True)
text(s, Inches(0.9), Inches(2.28), Inches(5.3), Inches(2.3), [
    [("data", {"bold": True, "color": GOLD}), ("  whitelist + vintage split", {})],
    [("woe / scorecard", {"bold": True, "color": GOLD}), ("  binning, IV screen, points", {})],
    [("models", {"bold": True, "color": GOLD}), ("  XGBoost, LightGBM, MLP", {})],
    [("calibration", {"bold": True, "color": GOLD}), ("  Platt + isotonic on 2014", {})],
    [("metrics", {"bold": True, "color": GOLD}), ("  AUC, KS, Brier, ECE, DeLong", {})],
    [("decision", {"bold": True, "color": GOLD}), ("  profit sweep over cutoffs", {})],
    [("pipeline", {"bold": True, "color": GOLD}), ("  orchestration, 5 seeds each", {})],
], size=12.5, space_after=5)
bullets(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(2.9), [
    ("The baseline is honest:", "an industry-style WoE scorecard, not a "
     "strawman logistic regression on raw inputs."),
    ("Same information for everyone:", "all four models see the identical 23 "
     "features; preprocessing cannot explain the gaps."),
    ("Calibration is separated from fitting", "and profit uses real cash flows: "
     "repaid loans earn their interest, charged-off lose 65% of principal."),
], size=13.5, space_after=9)

# the whitelist decision, made concrete
box(s, Inches(0.6), Inches(4.9), Inches(5.9), Inches(2.3), fill=LIGHT, round_=True)
text(s, Inches(0.9), Inches(5.05), Inches(5.3), Inches(0.35),
     "Kept: 23 origination-time fields", size=14.5, color=NAVY, bold=True)
text(s, Inches(0.9), Inches(5.45), Inches(5.3), Inches(1.6), [
    [("loan_amnt · term · int_rate · installment · grade · sub_grade", {})],
    [("annual_inc · dti · emp_length · home_ownership · purpose", {})],
    [("fico · revol_util · inq_last_6mths · delinq_2yrs · open_acc · "
      "revol_bal · total_acc · mort_acc · pub_rec", {})],
], size=11.5, color=INK, space_after=5)

box(s, Inches(6.9), Inches(4.9), Inches(5.8), Inches(2.3), fill=LIGHT, round_=True)
text(s, Inches(7.2), Inches(5.05), Inches(5.2), Inches(0.35),
     "Discarded: 128 columns", size=14.5, color=RED, bold=True)
text(s, Inches(7.2), Inches(5.45), Inches(5.2), Inches(1.6), [
    [("total_pymnt · total_rec_prncp · recoveries · "
      "collection_recovery_fee · last_pymnt_amnt · last_fico_range_high", {})],
    [("Known only after the loan was funded. ", {"italic": True}),
     ("total_rec_prncp", {"bold": True, "color": NAVY}),
     (" equals the loan amount for every repaid loan: it is the answer, not a "
      "feature.", {"italic": True})],
], size=11.5, color=INK, space_after=6)
notes(s, "Walk through the package on screen here. The design principles: an "
         "honest scorecard baseline, identical information for all four models, "
         "calibration separated from fitting, and profit from real cash flows. "
         "Then the concrete version of the leakage decision: on the left what I "
         "kept, all knowable on the day of application; on the right what I "
         "threw away. Point at total_rec_prncp: principal received to date "
         "equals the loan amount for every repaid loan, so a model using it "
         "scores near-perfectly and has learned nothing.")

# ---------------------------------------------------------------- 6: metrics
s = slide()
header(s, "Evaluation", "Three Questions, Each With Its Own Metrics")
for i, (t1, t2) in enumerate([
        ("Discrimination: AUC, KS",
         "AUC, the area under the receiver operating characteristic curve: the "
         "probability a random defaulter is ranked riskier than a random good "
         "borrower. KS, Kolmogorov-Smirnov: the maximum distance between the "
         "two score distributions. Both ignore whether the PD values "
         "themselves are right."),
        ("Calibration: Brier, ECE, reliability",
         "Among loans assigned PD near q, do a fraction q actually default? "
         "Brier: mean squared error of the probability. ECE: average gap "
         "between predicted PD and observed default rate over 10 bins."),
        ("Decision quality: profit, with significance",
         "Realized profit per 1,000 applicants across accept/reject cutoffs, "
         "from actual cash flows. Significance throughout: DeLong test for "
         "correlated AUCs, 1,000-resample bootstrap CIs.")]):
    y = Inches(1.9) + i * Inches(1.55)
    box(s, Inches(0.6), y, Inches(12.1), Inches(1.35), fill=LIGHT, round_=True)
    text(s, Inches(0.95), y + Inches(0.16), Inches(3.5), Inches(1.05), t1,
         size=16.5, color=NAVY, bold=True)
    text(s, Inches(4.6), y + Inches(0.14), Inches(7.7), Inches(1.1), t2,
         size=13.5, color=INK)
text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
     [[("Why three?  ", {"bold": True, "color": NAVY}),
       ("A model can gain AUC while its probabilities become less "
        "trustworthy; 'best model' is undefined until the question is fixed.",
        {})]], size=15)
notes(s, "Define the metrics before showing numbers: AUC and KS for ranking, "
         "Brier and ECE for calibration, realized profit for decisions, with "
         "DeLong and bootstrap significance everywhere. The point: the three "
         "questions can rank models differently.")

# ---------------------------------------------------------------- 7: leaderboard
s = slide()
header(s, "Results 1 of 4", "Discrimination: the Trees Win, the Neural Net Loses")
rows = [
    ("Scorecard", "0.6767", "0.2570", "baseline", INK),
    ("XGBoost", "0.6850", "0.2701", "+0.85 pts,  p < 10⁻⁵⁴", GREEN),
    ("LightGBM", "0.6843", "0.2693", "+0.72 pts,  p < 10⁻⁴⁸", GREEN),
    ("MLP", "0.6716", "0.2542", "−0.17 pts,  p = 0.024 (worse)", RED),
]
tbl_shape = s.shapes.add_table(5, 4, Inches(0.6), Inches(1.9), Inches(7.6), Inches(3.1))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.9)
tbl.columns[1].width = Inches(1.3)
tbl.columns[2].width = Inches(1.3)
tbl.columns[3].width = Inches(3.1)
hdr = ["Model", "AUC", "KS", "ΔAUC vs scorecard (DeLong)"]
for j, h in enumerate(hdr):
    c = tbl.cell(0, j)
    c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = WHITE
            r.font.name = "Calibri"
for i, (m, auc_, ks, d, col) in enumerate(rows):
    vals = [m, auc_, ks, d]
    for j, v in enumerate(vals):
        c = tbl.cell(i + 1, j)
        c.text = v
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14.5)
                r.font.name = "Calibri"
                r.font.bold = (j == 0)
                r.font.color.rgb = col if j == 3 else INK
bullets(s, Inches(8.6), Inches(2.0), Inches(4.2), Inches(4.3), [
    ("Statistically decisive:", "283,026 test loans give enormous power; the "
     "boosted-tree edge is unambiguous."),
    ("Economically modest:", "+0.85 AUC points, not the +8 to +15 sometimes "
     "quoted from random-split studies."),
    ("Deep learning does not pay:", "the MLP is significantly worse than the "
     "forty-year-old scorecard, on tabular credit data."),
], size=14.5, space_after=12)
text(s, Inches(0.6), Inches(5.4), Inches(7.6), Inches(1.2),
     [[("Why so close?  ", {"bold": True, "color": NAVY}),
       ("The features include the platform's own risk pricing (grade, interest "
        "rate), and the out-of-time split denies flexible models the flattery "
        "of a random split.", {})]], size=14.5)
notes(s, "First results: XGBoost +0.85 AUC points over the scorecard, LightGBM "
         "close behind, both overwhelmingly significant by DeLong. The MLP "
         "actually loses to the scorecard. Explain why the gap is modest.")

# ---------------------------------------------------------------- 8: calibration
s = slide()
header(s, "Results 2 of 4", "Calibration: Every Raw Model Under-Predicts the Future")
s.shapes.add_picture(str(FIGURES / "fig_reliability.png"), Inches(0.6), Inches(1.85),
                     height=Inches(4.35))
bullets(s, Inches(6.6), Inches(1.95), Inches(6.1), Inches(4.5), [
    ("Read the diagonal:", "a calibrated model's curve lies on it. Every raw "
     "curve (grey) sits above: the 2015 vintage defaulted more than the "
     "training years, so all models under-predict PD."),
    ("Raw trees are worst:", "ECE 0.034 vs the scorecard's 0.025; margin-"
     "maximizing ensembles are not probability models."),
    ("Isotonic recalibration helps everyone", "but cannot anticipate drift that "
     "was not yet visible in 2014."),
    ("The scorecard beats the trees on calibration:", "ECE 0.0175 vs XGBoost's 0.0214 "
     "after identical recalibration. The trees keep the better Brier score, "
     "which also rewards discrimination."),
], size=14.5, space_after=11)
notes(s, "Reliability diagrams: raw curves above the diagonal because of "
         "vintage drift. Isotonic recalibration on 2014 removes most but not "
         "all bias. The trees never match the scorecard on calibration; this is "
         "why regulators monitor rather than trust PD models.")

# ---------------------------------------------------------------- 9: profit
s = slide()
header(s, "Results 3 of 4", "Profit: the AUC Gap Becomes Money at the Cutoff")
s.shapes.add_picture(str(FIGURES / "fig_profit.png"), Inches(0.6), Inches(1.9),
                     width=Inches(8.3))
bullets(s, Inches(9.3), Inches(2.0), Inches(3.5), Inches(4.4), [
    ("Accepting everyone already earns", "$705 per applicant; interest income "
     "dominates this portfolio."),
    ("Scorecard:", "reject the riskiest 6.7% for +1.2% profit."),
    ("XGBoost:", "reject 9.3% for +3.3%, that is +$14,700 per 1,000 applicants "
     "over the scorecard (+2.1%)."),
    ("Scale decides adoption:", "a modest edge worth millions at bank volume, "
     "and nothing at credit-union volume."),
], size=14.5, space_after=11)
notes(s, "Profit curves from actual cash flows. The whole battle happens in the "
         "zoomed panel: XGBoost's optimal policy makes 2.1% more than the "
         "scorecard's. Modest per applicant, linear in volume.")

# ---------------------------------------------------------------- 10: interpret
s = slide()
header(s, "Results 4 of 4", "What Did the Models Learn? Mostly the Same Thing")
s.shapes.add_picture(str(FIGURES / "fig_importance.png"), Inches(0.6), Inches(1.9),
                     height=Inches(4.3))
bullets(s, Inches(8.6), Inches(2.0), Inches(4.2), Inches(4.4), [
    ("Both families agree", "that the platform's own pricing (interest rate, "
     "grade) carries the most signal; rank correlation 0.56."),
    ("They disagree in encoding:", "the scorecard leans on sub-grade and FICO "
     "bins; XGBoost prefers the continuous interest rate and income."),
    ("Interpretability trade-off:", "the scorecard is a readable sum of points "
     "per applicant; the ensemble needs post-hoc explanation tools."),
], size=14.5, space_after=12)
notes(s, "Feature importance vs information value: both models rely on the "
         "platform's risk pricing, so all four models largely re-rank loans "
         "Lending Club already ranked. The scorecard stays point-readable.")

# ---------------------------------------------------------------- 11: ablation
s = slide()
header(s, "Ablation", "What If We Hide Lending Club's Own Pricing?")
text(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.55),
     [[("The concern:  ", {"bold": True, "color": NAVY}),
       ("grade, sub-grade and interest rate are one variable in three "
        "notations, the platform's own risk assessment. They supply 36% of "
        "XGBoost's gain and 43% of LightGBM's. So are the models assessing "
        "credit risk, or just re-ranking loans Lending Club already ranked?",
        {})]], size=14.5)
for i, (t1, t2) in enumerate([
        ("Variant A", "Drop grade, sub-grade, interest rate.  20 features left."),
        ("Variant B", "Also drop installment: with the term fixed at 36 months "
         "it tracks the rate at ρ = 0.998, so A still leaks pricing "
         "through it. B is the honest test.")]):
    y = Inches(2.5) + i * Inches(1.15)
    box(s, Inches(0.6), y, Inches(5.8), Inches(0.98), fill=LIGHT, round_=True)
    text(s, Inches(0.9), y + Inches(0.13), Inches(1.5), Inches(0.4), t1,
         size=15, color=GOLD, bold=True)
    text(s, Inches(2.35), y + Inches(0.11), Inches(3.8), Inches(0.8), t2,
         size=12, color=INK)

rows = [("XGBoost", "0.6850", "0.6637", "−0.021", INK),
        ("LightGBM", "0.6843", "0.6622", "−0.022", INK),
        ("MLP", "0.6716", "0.6428", "−0.029", INK),
        ("Scorecard", "0.6767", "0.6397", "−0.037", RED)]
tbl = s.shapes.add_table(5, 4, Inches(6.9), Inches(2.5), Inches(5.8),
                         Inches(2.3)).table
for j, (h, w) in enumerate([("Model", 1.9), ("Baseline", 1.2),
                            ("Variant B", 1.3), ("Change", 1.4)]):
    tbl.columns[j].width = Inches(w)
    c = tbl.cell(0, j)
    c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(12.5)
            r.font.color.rgb = WHITE; r.font.name = "Calibri"
for i, (m, b, v, d, col) in enumerate(rows):
    for j, val in enumerate([m, b, v, d]):
        c = tbl.cell(i + 1, j)
        c.text = val
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12.5); r.font.name = "Calibri"
                r.font.bold = (j == 0)
                r.font.color.rgb = col if j == 3 else INK

bullets(s, Inches(0.6), Inches(5.05), Inches(12.1), Inches(2.1), [
    ("The pricing was hiding the ML advantage, not creating it:", "the "
     "scorecard degrades almost twice as fast as the trees, so XGBoost's lead "
     "grows from +0.83 to +2.39 AUC points, and its profit edge to $15,400 "
     "per 1,000 applicants."),
    ("The trees used the proxy; the scorecard could not:", "the scorecard "
     "scores 0.6397 under both variants because its information-value screen "
     "had already discarded installment as too weak. The old-fashioned screen "
     "was the protection."),
], size=14, space_after=9)
notes(s, "The honest follow-up experiment. The platform's grade and rate are "
         "legitimate inputs for an investor, but they narrow the question. "
         "Rerunning without them: every model gets worse, the scorecard worst, "
         "so the ML advantage nearly triples. And the scorecard's IV screen "
         "had already thrown out the installment proxy that both tree "
         "ensembles were quietly using.")

# ---------------------------------------------------------------- 12: robustness
s = slide()
header(s, "Robustness and limitations", "Same Ranking on an Independent Dataset")
box(s, Inches(0.6), Inches(1.85), Inches(5.9), Inches(2.6), fill=LIGHT, round_=True)
text(s, Inches(0.9), Inches(2.05), Inches(5.3), Inches(0.4),
     "UCI Taiwan credit-card default (30k clients)", size=16, color=NAVY, bold=True)
text(s, Inches(0.9), Inches(2.55), Inches(5.3), Inches(1.8), [
    [("XGBoost 0.788", {"bold": True, "color": GREEN}),
     ("  vs scorecard 0.772 AUC: +1.50 pts, p = 2.8×10⁻⁶", {})],
    [("MLP 0.766", {"bold": True, "color": RED}), ("  again below the scorecard", {})],
    [("Scorecard again best ECE", {"bold": True, "color": NAVY}),
     ("  after identical recalibration", {})],
], size=14.5, space_after=10)
text(s, Inches(6.9), Inches(1.5), Inches(5.9), Inches(0.4), "Honest limitations",
     size=16, color=NAVY, bold=True)
bullets(s, Inches(6.9), Inches(2.05), Inches(5.9), Inches(5.0), [
    ("Accepted loans only:", "all models inherit Lending Club's own selection; "
     "reject inference is an open problem."),
    ("Simple profit model:", "undiscounted, fixed LGD of 0.65, no partial "
     "recoveries or prepayment."),
    ("Platform grades kept as features:", "realistic for an investor; the "
     "ablation shows this understates the ML advantage by about 3x."),
    ("One market, one test vintage:", "2015 was benign; a recession vintage "
     "could reorder the calibration findings."),
    ("My scorecard is unconstrained:", "expert binning with monotonicity "
     "constraints would likely close part of the gap."),
], size=14, space_after=9)
notes(s, "Taiwan reproduces the ranking, so it is not a Lending Club artifact. "
         "Then the honest limitations: reject inference, simple cash-flow "
         "model, platform grades as features, single benign test vintage.")

# ---------------------------------------------------------------- 13: wrap-up
s = slide()
box(s, 0, 0, SW, SH, fill=NAVY)
text(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.35), "WRAP-UP",
     size=13, color=GOLD, bold=True)
text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.4),
     "A Qualified Yes: ML Wins, but Narrowly, and Only the Trees",
     size=31, color=WHITE, bold=True)
text(s, Inches(0.8), Inches(2.15), Inches(11.7), Inches(2.9), [
    [("Boosted trees beat the scorecard on all three criteria:  ", {"bold": True, "color": GOLD}),
     ("discrimination (+0.85 AUC pts, p < 10⁻⁵⁴), Brier after identical "
      "recalibration, and profit (+2.1% at the optimal cutoff).", {})],
    [("But the margins are modest and uneven:  ", {"bold": True, "color": GOLD}),
     ("the MLP loses to the scorecard outright, and the scorecard stays "
      "better calibrated than either tree ensemble.", {})],
    [("The margin depends on the information set:  ", {"bold": True, "color": GOLD}),
     ("withholding the platform's own pricing widens XGBoost's lead from "
      "0.83 to 2.39 AUC points, so the headline number understates what ML "
      "offers a lender underwriting from bureau data directly.", {})],
    [("What I learned:  ", {"bold": True, "color": GOLD}),
     ("leakage control was the highest-leverage work in the project; the model "
      "ranking changes with the question asked; and vintage drift is why PD "
      "models are monitored, not trusted.", {})],
    [("Future work:  ", {"bold": True, "color": GOLD}),
     ("reject inference, default-timing survival models, discounted cash "
      "flows, monotonic gradient boosting, a recession-vintage stress test.", {})],
], size=15.5, color=WHITE, space_after=13, line_spacing=1.08)
box(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.02), fill=GOLD)
text(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.8), [
    [("Code, results and report:  ", {"bold": True, "color": GOLD}),
     ("github.com/iitgF/T8", {"color": WHITE, "bold": True}),
     ("      François Schmitt  ·  23035010523  ·  BSc (Hons) DSAI, IIT Guwahati",
      {"color": RGBColor(0xC8, 0xD4, 0xE8), "size": 13})],
], size=15)
notes(s, "Wrap up: qualified yes. Recap the three criteria, the learnings, and "
         "future work. Point to the GitHub repo where everything reproduces "
         "with one command. Thank the viewer.")

prs.save(str(HERE / "Project_Video_Presentation.pptx"))
print("saved Project_Video_Presentation.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
